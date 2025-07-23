"""
Confocal Image Summary Generator (PowerPoint)
=============================================

This script creates a PowerPoint presentation (16:9) to summarize confocal microscopy images.
Each slide shows one sample, displaying four channels in a row: LUVs, Ld, BF, and merged.

Requirements:
-------------
- Python 3.x
- Package: python-pptx (`pip install python-pptx`)
- Images must be .png files and follow naming convention:
    basename_<suffix>.png
    where suffix ∈ ["_Atto488DPPE", "_Atto655DOPE", "_BF", "_merged"]

Usage:
------
1. Place your .png images in a single folder.
2. Modify `input_folder` to that folder path.
3. Run the script:

    python generate_confocal_summary_ppt.py

Output:
-------
A PowerPoint file named `Confocal_Summary.pptx` saved in the same folder.

"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt

# === Input image folder ===
input_folder = r"Path\To\Your\ImageFolder"  # ← Change this path before use

# === Channel suffixes and labels (in order) ===
channel_suffixes = ["_Atto488DPPE", "_Atto655DOPE", "_BF", "_merged"]
channel_labels = ["A", "B", "C", "D"]

# === Create PowerPoint (16:9) ===
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# === Group images by sample base name ===
sample_dict = {}

for file in os.listdir(input_folder):
    if file.lower().endswith(".png"):
        for suffix in channel_suffixes:
            if suffix in file:
                base = file.replace(suffix + ".png", "")
                sample_dict.setdefault(base, {})[suffix] = file
                break

# === Create slides ===
for sample_name, channels in sorted(sample_dict.items()):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide

    # Title: Sample name
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.5), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = f"Sample: {sample_name}"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = 1  # Center

    # Channel labels
    for i, label in enumerate(channel_labels):
        left = Inches(1.0 + i * 2.8)
        top = Inches(1.0)
        width = Inches(2.5)
        height = Inches(0.4)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = 1

    # Images
    for i, suffix in enumerate(channel_suffixes):
        if suffix in channels:
            img_path = os.path.join(input_folder, channels[suffix])
            left = Inches(1.0 + i * 2.8)
            top = Inches(1.5)
            width = Inches(2.5)
            slide.shapes.add_picture(img_path, left, top, width=width)

    # Notes section with filenames
    notes_slide = slide.notes_slide
    notes_text = "\n".join([f"{label}: {channels.get(suffix, 'MISSING')}" 
                            for suffix, label in zip(channel_suffixes, channel_labels)])
    notes_slide.notes_text_frame.text = f"Sample: {sample_name}\n{notes_text}"

# === Save PowerPoint file ===
output_path = os.path.join(input_folder, "Confocal_Summary.pptx")
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
